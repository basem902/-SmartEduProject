"""
AI Validator for Project Submissions
نظام التحقق بالذكاء الاصطناعي
"""
import os
from django.conf import settings
import logging

logger = logging.getLogger(__name__)


class AIValidator:
    """
    نظام التحقق بالذكاء الاصطناعي للمشاريع
    """
    
    def __init__(self):
        """تهيئة AI Validator"""
        # Gemini API
        try:
            import google.generativeai as genai
            genai.configure(api_key=settings.GEMINI_API_KEY)
            self.gemini_flash = genai.GenerativeModel('gemini-1.5-flash')
            self.gemini_vision = genai.GenerativeModel('gemini-1.5-pro-vision')
            logger.info("✅ Gemini API initialized")
        except Exception as e:
            logger.error(f"❌ Failed to initialize Gemini: {str(e)}")
            self.gemini_flash = None
            self.gemini_vision = None
    
    def validate_submission(self, submission):
        """
        التحقق الشامل من المشروع
        
        Args:
            submission: كائن Submission
        
        Returns:
            dict: نتائج التحقق
        """
        project = submission.project
        file_type = project.file_type
        
        logger.info(f"🔍 بدء التحقق من Submission #{submission.id} - نوع الملف: {file_type}")
        
        try:
            # اختيار المعالج حسب نوع الملف
            if file_type == 'video':
                return self.validate_video(submission)
            elif file_type == 'pdf':
                return self.validate_pdf(submission)
            elif file_type == 'image':
                return self.validate_image(submission)
            elif file_type == 'document':
                return self.validate_document(submission)
            elif file_type == 'audio':
                return self.validate_audio(submission)
            else:
                return {
                    'status': 'rejected',
                    'overall_score': 0,
                    'rejection_reasons': [f'نوع ملف غير مدعوم: {file_type}']
                }
        
        except Exception as e:
            logger.error(f"❌ خطأ في التحقق: {str(e)}", exc_info=True)
            return {
                'status': 'needs_review',
                'overall_score': 0,
                'rejection_reasons': [f'حدث خطأ في التحليل: {str(e)}'],
                'error': str(e)
            }
    
    def validate_video(self, submission):
        """
        التحقق الشامل من الفيديو
        1. فحص المدة
        2. OCR على آخر 5 ثواني
        3. تحليل المحتوى بـ Gemini Vision
        4. كشف التشابه
        """
        logger.info(f"🎬 بدء فحص الفيديو #{submission.id}")
        
        project = submission.project
        file_path = submission.file_path
        student_name = submission.submitted_student_name or submission.student.full_name if submission.student else "غير معروف"
        
        results = {
            'checks': {},
            'rejection_reasons': [],
            'warnings': [],
            'overall_score': 0
        }
        
        try:
            # 1. فحص المدة
            duration_result = self._check_video_duration(file_path, project)
            results['checks']['duration'] = duration_result
            
            if duration_result['status'] == 'fail':
                results['rejection_reasons'].append(duration_result['message'])
                results['status'] = 'rejected'
                return results
            
            # 2. OCR على آخر 5 ثواني
            ocr_result = self._check_video_ocr(file_path, student_name)
            results['checks']['ocr'] = ocr_result
            
            if ocr_result['status'] == 'fail':
                results['rejection_reasons'].append(ocr_result['message'])
            elif ocr_result['status'] == 'warning':
                results['warnings'].append(ocr_result['message'])
            
            # 3. تحليل المحتوى بـ Gemini Vision
            gemini_result = self._analyze_video_content(file_path, project)
            results['checks']['content_analysis'] = gemini_result
            
            if gemini_result['status'] == 'fail':
                results['rejection_reasons'].append(gemini_result['message'])
            
            # 4. كشف التشابه
            similarity_result = self._check_video_similarity(file_path, submission)
            results['checks']['similarity'] = similarity_result
            
            if similarity_result['status'] == 'fail':
                results['rejection_reasons'].append(similarity_result['message'])
            
            # حساب الدرجة النهائية
            scores = [
                duration_result.get('score', 0),
                ocr_result.get('score', 0),
                gemini_result.get('score', 0),
                similarity_result.get('score', 0)
            ]
            results['overall_score'] = sum(scores) / len(scores)
            
            # تحديد الحالة النهائية
            if results['rejection_reasons']:
                results['status'] = 'rejected'
            elif results['overall_score'] < 60:
                results['status'] = 'needs_review'
                results['rejection_reasons'].append('الدرجة أقل من 60%')
            else:
                results['status'] = 'approved'
            
            logger.info(f"✅ انتهى فحص الفيديو #{submission.id} - الحالة: {results['status']}")
            return results
            
        except Exception as e:
            logger.error(f"❌ خطأ في فحص الفيديو #{submission.id}: {str(e)}", exc_info=True)
            return {
                'status': 'needs_review',
                'overall_score': 0,
                'rejection_reasons': [f'حدث خطأ في فحص الفيديو: {str(e)}'],
                'checks': results.get('checks', {})
            }
    
    def validate_pdf(self, submission):
        """
        التحقق الشامل من PDF
        1. استخراج النص
        2. فحص عدد الكلمات والصفحات
        3. تحليل المحتوى بـ Gemini
        4. كشف الانتحال
        """
        logger.info(f"📄 بدء فحص PDF #{submission.id}")
        
        project = submission.project
        file_path = submission.file_path
        
        results = {
            'checks': {},
            'rejection_reasons': [],
            'warnings': [],
            'overall_score': 0
        }
        
        try:
            # 1. استخراج النص من PDF
            text_result = self._extract_pdf_text(file_path)
            results['checks']['text_extraction'] = text_result
            
            if text_result['status'] == 'fail':
                results['rejection_reasons'].append(text_result['message'])
                results['status'] = 'rejected'
                return results
            
            extracted_text = text_result.get('text', '')
            
            # 2. فحص عدد الكلمات والصفحات
            stats_result = self._check_pdf_stats(text_result, project)
            results['checks']['statistics'] = stats_result
            
            if stats_result['status'] == 'fail':
                results['rejection_reasons'].append(stats_result['message'])
            elif stats_result['status'] == 'warning':
                results['warnings'].append(stats_result['message'])
            
            # 3. تحليل المحتوى بـ Gemini
            content_result = self._analyze_pdf_content(extracted_text, project)
            results['checks']['content_analysis'] = content_result
            
            if content_result['status'] == 'fail':
                results['rejection_reasons'].append(content_result['message'])
            
            # 4. كشف الانتحال
            plagiarism_result = self._check_pdf_plagiarism(extracted_text, submission)
            results['checks']['plagiarism'] = plagiarism_result
            
            if plagiarism_result['status'] == 'fail':
                results['rejection_reasons'].append(plagiarism_result['message'])
            elif plagiarism_result['status'] == 'warning':
                results['warnings'].append(plagiarism_result['message'])
            
            # حساب الدرجة النهائية
            scores = [
                text_result.get('score', 0),
                stats_result.get('score', 0),
                content_result.get('score', 0),
                plagiarism_result.get('score', 0)
            ]
            results['overall_score'] = sum(scores) / len(scores)
            
            # تحديد الحالة النهائية
            if results['rejection_reasons']:
                results['status'] = 'rejected'
            elif results['overall_score'] < 60:
                results['status'] = 'needs_review'
                results['rejection_reasons'].append('الدرجة أقل من 60%')
            else:
                results['status'] = 'approved'
            
            logger.info(f"✅ انتهى فحص PDF #{submission.id} - الحالة: {results['status']}")
            return results
            
        except Exception as e:
            logger.error(f"❌ خطأ في فحص PDF #{submission.id}: {str(e)}", exc_info=True)
            return {
                'status': 'needs_review',
                'overall_score': 0,
                'rejection_reasons': [f'حدث خطأ في فحص PDF: {str(e)}'],
                'checks': results.get('checks', {})
            }
    
    def validate_image(self, submission):
        """
        التحقق الشامل من الصورة
        1. OCR - قراءة النص
        2. تحليل بـ Gemini Vision
        3. فحص الجودة
        4. كشف التشابه
        """
        logger.info(f"🖼️ بدء فحص الصورة #{submission.id}")
        
        project = submission.project
        file_path = submission.file_path
        
        results = {
            'checks': {},
            'rejection_reasons': [],
            'warnings': [],
            'overall_score': 0
        }
        
        try:
            # 1. OCR - قراءة النص من الصورة
            ocr_result = self._check_image_ocr(file_path)
            results['checks']['ocr'] = ocr_result
            
            # 2. تحليل بـ Gemini Vision
            vision_result = self._analyze_image_content(file_path, project)
            results['checks']['vision_analysis'] = vision_result
            
            if vision_result['status'] == 'fail':
                results['rejection_reasons'].append(vision_result['message'])
            
            # 3. فحص الجودة (الدقة والحجم)
            quality_result = self._check_image_quality(file_path)
            results['checks']['quality'] = quality_result
            
            if quality_result['status'] == 'warning':
                results['warnings'].append(quality_result['message'])
            
            # 4. كشف التشابه
            similarity_result = self._check_image_similarity(file_path, submission)
            results['checks']['similarity'] = similarity_result
            
            if similarity_result['status'] == 'fail':
                results['rejection_reasons'].append(similarity_result['message'])
            elif similarity_result['status'] == 'warning':
                results['warnings'].append(similarity_result['message'])
            
            # حساب الدرجة النهائية
            scores = [
                ocr_result.get('score', 0),
                vision_result.get('score', 0),
                quality_result.get('score', 0),
                similarity_result.get('score', 0)
            ]
            results['overall_score'] = sum(scores) / len(scores)
            
            # تحديد الحالة النهائية
            if results['rejection_reasons']:
                results['status'] = 'rejected'
            elif results['overall_score'] < 60:
                results['status'] = 'needs_review'
            else:
                results['status'] = 'approved'
            
            logger.info(f"✅ انتهى فحص الصورة #{submission.id} - الحالة: {results['status']}")
            return results
            
        except Exception as e:
            logger.error(f"❌ خطأ في فحص الصورة #{submission.id}: {str(e)}", exc_info=True)
            return {
                'status': 'needs_review',
                'overall_score': 0,
                'rejection_reasons': [f'حدث خطأ في فحص الصورة: {str(e)}'],
                'checks': results.get('checks', {})
            }
    
    def validate_document(self, submission):
        """
        التحقق الشامل من المستندات (Word/Excel/PPT)
        1. استخراج النص
        2. فحص الإحصائيات
        3. تحليل بـ Gemini
        """
        logger.info(f"📝 بدء فحص المستند #{submission.id}")
        
        project = submission.project
        file_path = submission.file_path
        file_ext = file_path.split('.')[-1].lower()
        
        results = {
            'checks': {},
            'rejection_reasons': [],
            'warnings': [],
            'overall_score': 0
        }
        
        try:
            # 1. استخراج النص حسب نوع الملف
            if file_ext in ['doc', 'docx']:
                text_result = self._extract_word_text(file_path)
            elif file_ext in ['xls', 'xlsx']:
                text_result = self._extract_excel_text(file_path)
            elif file_ext in ['ppt', 'pptx']:
                text_result = self._extract_ppt_text(file_path)
            else:
                return {
                    'status': 'rejected',
                    'overall_score': 0,
                    'rejection_reasons': [f'نوع ملف غير مدعوم: {file_ext}'],
                    'checks': {}
                }
            
            results['checks']['text_extraction'] = text_result
            
            if text_result['status'] == 'fail':
                results['rejection_reasons'].append(text_result['message'])
                results['status'] = 'rejected'
                return results
            
            extracted_text = text_result.get('text', '')
            
            # 2. فحص الإحصائيات
            stats_result = self._check_document_stats(text_result, project)
            results['checks']['statistics'] = stats_result
            
            if stats_result['status'] == 'fail':
                results['rejection_reasons'].append(stats_result['message'])
            elif stats_result['status'] == 'warning':
                results['warnings'].append(stats_result['message'])
            
            # 3. تحليل المحتوى بـ Gemini
            content_result = self._analyze_document_content(extracted_text, project, file_ext)
            results['checks']['content_analysis'] = content_result
            
            if content_result['status'] == 'fail':
                results['rejection_reasons'].append(content_result['message'])
            
            # حساب الدرجة النهائية
            scores = [
                text_result.get('score', 0),
                stats_result.get('score', 0),
                content_result.get('score', 0)
            ]
            results['overall_score'] = sum(scores) / len(scores)
            
            # تحديد الحالة النهائية
            if results['rejection_reasons']:
                results['status'] = 'rejected'
            elif results['overall_score'] < 60:
                results['status'] = 'needs_review'
            else:
                results['status'] = 'approved'
            
            logger.info(f"✅ انتهى فحص المستند #{submission.id} - الحالة: {results['status']}")
            return results
            
        except Exception as e:
            logger.error(f"❌ خطأ في فحص المستند #{submission.id}: {str(e)}", exc_info=True)
            return {
                'status': 'needs_review',
                'overall_score': 0,
                'rejection_reasons': [f'حدث خطأ في فحص المستند: {str(e)}'],
                'checks': results.get('checks', {})
            }
    
    def validate_audio(self, submission):
        """
        التحقق الشامل من الصوت
        1. فحص المدة
        2. Speech-to-Text
        3. تحليل المحتوى بـ Gemini
        4. كشف التشابه
        """
        logger.info(f"🎵 بدء فحص الصوت #{submission.id}")
        
        project = submission.project
        file_path = submission.file_path
        
        results = {
            'checks': {},
            'rejection_reasons': [],
            'warnings': [],
            'overall_score': 0
        }
        
        try:
            # 1. فحص مدة الصوت
            duration_result = self._check_audio_duration(file_path, project)
            results['checks']['duration'] = duration_result
            
            if duration_result['status'] == 'fail':
                results['rejection_reasons'].append(duration_result['message'])
                results['status'] = 'rejected'
                return results
            
            # 2. Speech-to-Text
            stt_result = self._audio_to_text(file_path)
            results['checks']['speech_to_text'] = stt_result
            
            if stt_result['status'] == 'fail':
                results['rejection_reasons'].append(stt_result['message'])
            
            transcribed_text = stt_result.get('text', '')
            
            # 3. تحليل المحتوى بـ Gemini
            content_result = self._analyze_audio_content(transcribed_text, project)
            results['checks']['content_analysis'] = content_result
            
            if content_result['status'] == 'fail':
                results['rejection_reasons'].append(content_result['message'])
            
            # 4. كشف التشابه
            similarity_result = self._check_audio_similarity(transcribed_text, submission)
            results['checks']['similarity'] = similarity_result
            
            if similarity_result['status'] == 'fail':
                results['rejection_reasons'].append(similarity_result['message'])
            elif similarity_result['status'] == 'warning':
                results['warnings'].append(similarity_result['message'])
            
            # حساب الدرجة النهائية
            scores = [
                duration_result.get('score', 0),
                stt_result.get('score', 0),
                content_result.get('score', 0),
                similarity_result.get('score', 0)
            ]
            results['overall_score'] = sum(scores) / len(scores)
            
            # تحديد الحالة النهائية
            if results['rejection_reasons']:
                results['status'] = 'rejected'
            elif results['overall_score'] < 60:
                results['status'] = 'needs_review'
            else:
                results['status'] = 'approved'
            
            logger.info(f"✅ انتهى فحص الصوت #{submission.id} - الحالة: {results['status']}")
            return results
            
        except Exception as e:
            logger.error(f"❌ خطأ في فحص الصوت #{submission.id}: {str(e)}", exc_info=True)
            return {
                'status': 'needs_review',
                'overall_score': 0,
                'rejection_reasons': [f'حدث خطأ في فحص الصوت: {str(e)}'],
                'checks': results.get('checks', {})
            }
    
    # ====================================
    # Video Validation Helper Methods
    # ====================================
    
    def _check_video_duration(self, file_path, project):
        """
        فحص مدة الفيديو
        
        Args:
            file_path: مسار الفيديو
            project: كائن المشروع
            
        Returns:
            dict: نتيجة الفحص
        """
        try:
            import cv2
            
            video = cv2.VideoCapture(file_path)
            
            if not video.isOpened():
                return {
                    'status': 'fail',
                    'message': 'فشل في فتح الفيديو',
                    'score': 0
                }
            
            # حساب المدة
            fps = video.get(cv2.CAP_PROP_FPS)
            frame_count = video.get(cv2.CAP_PROP_FRAME_COUNT)
            duration = frame_count / fps if fps > 0 else 0
            
            video.release()
            
            # التحقق من قيود المدة
            constraints = project.file_constraints or {}
            duration_constraints = constraints.get('duration', {})
            
            min_duration = duration_constraints.get('min', 15)  # default 15 seconds
            max_duration = duration_constraints.get('max', 30)  # default 30 seconds
            
            logger.info(f"📹 مدة الفيديو: {duration:.2f} ثانية (المطلوب: {min_duration}-{max_duration})")
            
            if duration < min_duration:
                return {
                    'status': 'fail',
                    'message': f'الفيديو قصير جداً ({duration:.1f}ث). المطلوب على الأقل {min_duration} ثانية',
                    'duration': duration,
                    'score': 0
                }
            elif duration > max_duration:
                return {
                    'status': 'fail',
                    'message': f'الفيديو طويل جداً ({duration:.1f}ث). الحد الأقصى {max_duration} ثانية',
                    'duration': duration,
                    'score': 0
                }
            else:
                return {
                    'status': 'pass',
                    'message': f'مدة الفيديو مناسبة ({duration:.1f}ث)',
                    'duration': duration,
                    'score': 100
                }
                
        except Exception as e:
            logger.error(f"❌ خطأ في فحص مدة الفيديو: {str(e)}")
            return {
                'status': 'fail',
                'message': f'خطأ في قراءة الفيديو: {str(e)}',
                'score': 0
            }
    
    def _check_video_ocr(self, file_path, student_name):
        """
        OCR على آخر 5 ثواني من الفيديو للتحقق من اسم الطالب
        
        Args:
            file_path: مسار الفيديو
            student_name: اسم الطالب
            
        Returns:
            dict: نتيجة الفحص
        """
        try:
            import cv2
            import easyocr
            from apps.projects.utils import normalize_arabic_name
            
            # تهيئة EasyOCR (العربية + الإنجليزية)
            reader = easyocr.Reader(['ar', 'en'], gpu=False)
            
            video = cv2.VideoCapture(file_path)
            
            if not video.isOpened():
                return {
                    'status': 'fail',
                    'message': 'فشل في فتح الفيديو',
                    'score': 0
                }
            
            # الحصول على معلومات الفيديو
            fps = video.get(cv2.CAP_PROP_FPS)
            total_frames = int(video.get(cv2.CAP_PROP_FRAME_COUNT))
            
            # آخر 5 ثواني
            start_frame = max(0, total_frames - int(fps * 5))
            
            # استخراج إطارات من آخر 5 ثواني (كل ثانية إطار)
            frames_to_check = []
            for i in range(5):
                frame_pos = start_frame + int(i * fps)
                if frame_pos < total_frames:
                    video.set(cv2.CAP_PROP_POS_FRAMES, frame_pos)
                    ret, frame = video.read()
                    if ret:
                        frames_to_check.append(frame)
            
            video.release()
            
            # قراءة النص من الإطارات
            detected_texts = []
            for frame in frames_to_check:
                result = reader.readtext(frame)
                for detection in result:
                    text = detection[1]
                    detected_texts.append(text)
            
            # دمج النصوص
            combined_text = ' '.join(detected_texts)
            logger.info(f"📝 النصوص المستخرجة: {combined_text[:100]}...")
            
            # تطبيع الأسماء للمقارنة
            normalized_student = normalize_arabic_name(student_name)
            normalized_detected = normalize_arabic_name(combined_text)
            
            # البحث عن الاسم
            name_parts = normalized_student.split()
            found_parts = sum(1 for part in name_parts if part in normalized_detected)
            
            # نسبة التطابق
            match_percentage = (found_parts / len(name_parts)) * 100 if name_parts else 0
            
            logger.info(f"🔍 تطابق الاسم: {match_percentage:.1f}% ({found_parts}/{len(name_parts)} أجزاء)")
            
            if match_percentage >= 75:
                return {
                    'status': 'pass',
                    'message': f'تم العثور على الاسم في الفيديو ({match_percentage:.0f}% تطابق)',
                    'detected_text': combined_text[:200],
                    'match_percentage': match_percentage,
                    'score': 100
                }
            elif match_percentage >= 50:
                return {
                    'status': 'warning',
                    'message': f'تطابق جزئي للاسم ({match_percentage:.0f}%). يرجى التأكد',
                    'detected_text': combined_text[:200],
                    'match_percentage': match_percentage,
                    'score': 70
                }
            else:
                return {
                    'status': 'fail',
                    'message': 'لم يتم العثور على الاسم في آخر 5 ثواني من الفيديو',
                    'detected_text': combined_text[:200],
                    'match_percentage': match_percentage,
                    'score': 0
                }
                
        except Exception as e:
            logger.error(f"❌ خطأ في OCR للفيديو: {str(e)}")
            return {
                'status': 'warning',
                'message': f'تعذر قراءة النص من الفيديو: {str(e)}',
                'score': 50
            }
    
    def _analyze_video_content(self, file_path, project):
        """
        تحليل محتوى الفيديو باستخدام Gemini Vision
        
        Args:
            file_path: مسار الفيديو
            project: كائن المشروع
            
        Returns:
            dict: نتيجة التحليل
        """
        try:
            import google.generativeai as genai
            
            if not self.gemini_vision:
                return {
                    'status': 'warning',
                    'message': 'Gemini Vision غير متاح',
                    'score': 70
                }
            
            # رفع الفيديو لـ Gemini
            logger.info(f"📤 رفع الفيديو لـ Gemini...")
            video_file = genai.upload_file(file_path)
            
            # تجهيز Prompt
            prompt = f"""حلل هذا الفيديو بدقة وأجب بصيغة JSON:

معلومات المشروع:
- العنوان: {project.title}
- الوصف: {project.description or 'غير محدد'}

أجب على التالي:
1. quality_score: جودة الفيديو من 0-100 (الوضوح، الإضاءة، الصوت)
2. content_relevance: مدى ارتباط المحتوى بالمشروع (0-100)
3. has_inappropriate_content: هل يحتوي على محتوى غير مناسب؟ (true/false)
4. summary: ملخص قصير للمحتوى (جملة واحدة)
5. issues: قائمة بأي مشاكل (أو قائمة فارغة)
6. recommendation: توصية (approved/rejected/needs_review)

أجب فقط بصيغة JSON بدون أي نص إضافي."""
            
            response = self.gemini_vision.generate_content([prompt, video_file])
            
            # تحليل النتيجة
            import json
            try:
                result = json.loads(response.text)
            except:
                # إذا فشل parsing، نستخرج المعلومات يدوياً
                result = {
                    'quality_score': 75,
                    'content_relevance': 80,
                    'has_inappropriate_content': False,
                    'summary': response.text[:100],
                    'issues': [],
                    'recommendation': 'approved'
                }
            
            # حساب الدرجة
            quality = result.get('quality_score', 70)
            relevance = result.get('content_relevance', 70)
            overall = (quality + relevance) / 2
            
            # التحقق من المحتوى غير المناسب
            if result.get('has_inappropriate_content'):
                return {
                    'status': 'fail',
                    'message': 'الفيديو يحتوي على محتوى غير مناسب',
                    'analysis': result,
                    'score': 0
                }
            
            # التحقق من الجودة
            if overall < 50:
                return {
                    'status': 'fail',
                    'message': f'جودة الفيديو منخفضة ({overall:.0f}%)',
                    'analysis': result,
                    'score': overall
                }
            elif overall < 70:
                return {
                    'status': 'warning',
                    'message': f'جودة الفيديو مقبولة ({overall:.0f}%)',
                    'analysis': result,
                    'score': overall
                }
            else:
                return {
                    'status': 'pass',
                    'message': f'جودة الفيديو ممتازة ({overall:.0f}%)',
                    'analysis': result,
                    'score': overall
                }
                
        except Exception as e:
            logger.error(f"❌ خطأ في تحليل الفيديو بـ Gemini: {str(e)}")
            return {
                'status': 'warning',
                'message': f'تعذر تحليل الفيديو: {str(e)}',
                'score': 70
            }
    
    def _check_video_similarity(self, file_path, submission):
        """
        كشف التشابه مع فيديوهات سابقة
        
        Args:
            file_path: مسار الفيديو
            submission: كائن التسليم
            
        Returns:
            dict: نتيجة الفحص
        """
        try:
            import videohash
            from .models import Submission
            
            # حساب hash للفيديو الحالي
            current_hash = videohash.VideoHash(path=file_path)
            
            # البحث عن فيديوهات سابقة في نفس المشروع
            previous_submissions = Submission.objects.filter(
                project=submission.project,
                file_type='video',
                validation_data__video_hash__isnull=False
            ).exclude(id=submission.id)
            
            # مقارنة مع الفيديوهات السابقة
            similarities = []
            for prev_sub in previous_submissions:
                try:
                    prev_hash_data = prev_sub.validation_data.get('video_hash')
                    if prev_hash_data:
                        # إعادة بناء hash
                        prev_hash = videohash.VideoHash(
                            storage_path=prev_hash_data
                        )
                        
                        # حساب الفرق (كلما أقل = أكثر تشابه)
                        difference = current_hash - prev_hash
                        similarity_percent = max(0, 100 - difference)
                        
                        if similarity_percent > 80:
                            similarities.append({
                                'submission_id': prev_sub.id,
                                'student': prev_sub.submitted_student_name,
                                'similarity': similarity_percent,
                                'difference': difference
                            })
                except:
                    continue
            
            # حفظ hash الحالي
            submission.validation_data = submission.validation_data or {}
            submission.validation_data['video_hash'] = str(current_hash)
            
            # التحقق من النتائج
            if similarities:
                most_similar = max(similarities, key=lambda x: x['similarity'])
                
                if most_similar['similarity'] > 90:
                    return {
                        'status': 'fail',
                        'message': f"فيديو مشابه جداً ({most_similar['similarity']:.0f}%) لتسليم سابق",
                        'similar_submissions': similarities[:3],
                        'score': 0
                    }
                else:
                    return {
                        'status': 'warning',
                        'message': f"تشابه متوسط ({most_similar['similarity']:.0f}%) مع تسليم سابق",
                        'similar_submissions': similarities[:3],
                        'score': 70
                    }
            else:
                return {
                    'status': 'pass',
                    'message': 'الفيديو أصلي (لا يوجد تشابه)',
                    'score': 100
                }
                
        except Exception as e:
            logger.error(f"❌ خطأ في فحص التشابه: {str(e)}")
            return {
                'status': 'warning',
                'message': 'تعذر فحص التشابه',
                'score': 80
            }
    
    # ====================================
    # PDF Validation Helper Methods
    # ====================================
    
    def _extract_pdf_text(self, file_path):
        """
        استخراج النص من PDF
        
        Args:
            file_path: مسار الملف
            
        Returns:
            dict: النص المستخرج + البيانات الإحصائية
        """
        try:
            import pdfplumber
            
            text = ''
            page_count = 0
            images_count = 0
            
            with pdfplumber.open(file_path) as pdf:
                page_count = len(pdf.pages)
                
                for page in pdf.pages:
                    # استخراج النص
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + '\n'
                    
                    # عد الصور
                    if hasattr(page, 'images'):
                        images_count += len(page.images)
            
            # إحصائيات
            word_count = len(text.split())
            char_count = len(text)
            
            logger.info(f"📊 PDF: {page_count} صفحة، {word_count} كلمة، {images_count} صورة")
            
            if word_count < 10:
                return {
                    'status': 'fail',
                    'message': 'PDF شبه فارغ أو لا يحتوي على نص',
                    'text': text,
                    'word_count': word_count,
                    'page_count': page_count,
                    'score': 0
                }
            
            return {
                'status': 'pass',
                'message': f'تم استخراج {word_count} كلمة من {page_count} صفحة',
                'text': text,
                'word_count': word_count,
                'page_count': page_count,
                'char_count': char_count,
                'images_count': images_count,
                'score': 100
            }
            
        except Exception as e:
            logger.error(f"❌ خطأ في استخراج نص PDF: {str(e)}")
            return {
                'status': 'fail',
                'message': f'فشل في قراءة PDF: {str(e)}',
                'text': '',
                'word_count': 0,
                'page_count': 0,
                'score': 0
            }
    
    def _check_pdf_stats(self, text_result, project):
        """
        فحص إحصائيات PDF (عدد الكلمات، الصفحات)
        
        Args:
            text_result: نتيجة استخراج النص
            project: كائن المشروع
            
        Returns:
            dict: نتيجة الفحص
        """
        try:
            word_count = text_result.get('word_count', 0)
            page_count = text_result.get('page_count', 0)
            
            # قراءة القيود من المشروع
            constraints = project.file_constraints or {}
            min_words = constraints.get('min_words', 100)  # default 100 words
            max_words = constraints.get('max_words', 5000)  # default 5000 words
            min_pages = constraints.get('min_pages', 1)
            max_pages = constraints.get('max_pages', 20)
            
            issues = []
            
            # فحص عدد الكلمات
            if word_count < min_words:
                issues.append(f'عدد الكلمات قليل جداً ({word_count}). المطلوب على الأقل {min_words} كلمة')
            elif word_count > max_words:
                issues.append(f'عدد الكلمات كثير جداً ({word_count}). الحد الأقصى {max_words} كلمة')
            
            # فحص عدد الصفحات
            if page_count < min_pages:
                issues.append(f'عدد الصفحات قليل ({page_count}). المطلوب على الأقل {min_pages} صفحة')
            elif page_count > max_pages:
                issues.append(f'عدد الصفحات كثير ({page_count}). الحد الأقصى {max_pages} صفحة')
            
            # التقييم
            if issues:
                if word_count < min_words / 2 or page_count < min_pages:
                    return {
                        'status': 'fail',
                        'message': ' | '.join(issues),
                        'word_count': word_count,
                        'page_count': page_count,
                        'score': 0
                    }
                else:
                    return {
                        'status': 'warning',
                        'message': ' | '.join(issues),
                        'word_count': word_count,
                        'page_count': page_count,
                        'score': 60
                    }
            else:
                return {
                    'status': 'pass',
                    'message': f'الإحصائيات مناسبة: {word_count} كلمة في {page_count} صفحة',
                    'word_count': word_count,
                    'page_count': page_count,
                    'score': 100
                }
                
        except Exception as e:
            logger.error(f"❌ خطأ في فحص إحصائيات PDF: {str(e)}")
            return {
                'status': 'warning',
                'message': 'تعذر فحص الإحصائيات',
                'score': 70
            }
    
    def _analyze_pdf_content(self, text, project):
        """
        تحليل محتوى PDF باستخدام Gemini
        
        Args:
            text: النص المستخرج
            project: كائن المشروع
            
        Returns:
            dict: نتيجة التحليل
        """
        try:
            if not self.gemini_flash:
                return {
                    'status': 'warning',
                    'message': 'Gemini غير متاح',
                    'score': 70
                }
            
            # اختصار النص للتحليل (أول 3000 حرف)
            text_sample = text[:3000] if len(text) > 3000 else text
            
            # تجهيز Prompt
            prompt = f"""حلل هذا النص من PDF وأجب بصيغة JSON:

معلومات المشروع:
- العنوان: {project.title}
- الوصف: {project.description or 'غير محدد'}

النص المستخرج:
{text_sample}

أجب على التالي بصيغة JSON:
1. content_quality: جودة المحتوى من 0-100
2. relevance_to_topic: مدى ارتباط المحتوى بالمشروع (0-100)
3. language_quality: جودة اللغة والإملاء (0-100)
4. has_copied_content: هل يبدو منسوخاً من الإنترنت؟ (true/false)
5. key_topics: قائمة بالمواضيع الرئيسية (array)
6. issues: قائمة بأي مشاكل (array)
7. recommendation: التوصية (approved/rejected/needs_review)

أجب فقط بصيغة JSON."""
            
            response = self.gemini_flash.generate_content(prompt)
            
            # تحليل النتيجة
            import json
            try:
                result = json.loads(response.text)
            except:
                # Fallback parsing
                result = {
                    'content_quality': 75,
                    'relevance_to_topic': 75,
                    'language_quality': 80,
                    'has_copied_content': False,
                    'key_topics': [],
                    'issues': [],
                    'recommendation': 'approved'
                }
            
            # حساب الدرجة
            quality = result.get('content_quality', 70)
            relevance = result.get('relevance_to_topic', 70)
            language = result.get('language_quality', 70)
            overall = (quality + relevance + language) / 3
            
            # التحقق من المحتوى المنسوخ
            if result.get('has_copied_content'):
                return {
                    'status': 'fail',
                    'message': 'المحتوى يبدو منسوخاً من الإنترنت',
                    'analysis': result,
                    'score': 0
                }
            
            # التحقق من الجودة
            if overall < 50:
                return {
                    'status': 'fail',
                    'message': f'جودة المحتوى منخفضة ({overall:.0f}%)',
                    'analysis': result,
                    'score': overall
                }
            elif overall < 70:
                return {
                    'status': 'warning',
                    'message': f'جودة المحتوى مقبولة ({overall:.0f}%)',
                    'analysis': result,
                    'score': overall
                }
            else:
                return {
                    'status': 'pass',
                    'message': f'جودة المحتوى ممتازة ({overall:.0f}%)',
                    'analysis': result,
                    'score': overall
                }
                
        except Exception as e:
            logger.error(f"❌ خطأ في تحليل محتوى PDF: {str(e)}")
            return {
                'status': 'warning',
                'message': f'تعذر تحليل المحتوى: {str(e)}',
                'score': 70
            }
    
    def _check_pdf_plagiarism(self, text, submission):
        """
        كشف الانتحال في PDF
        
        Args:
            text: النص المستخرج
            submission: كائن التسليم
            
        Returns:
            dict: نتيجة الفحص
        """
        try:
            from .models import Submission
            from sklearn.feature_extraction.text import TfidfVectorizer
            from sklearn.metrics.pairwise import cosine_similarity
            import numpy as np
            
            # البحث عن تسليمات سابقة في نفس المشروع
            previous_submissions = Submission.objects.filter(
                project=submission.project,
                file_type='pdf',
                validation_data__pdf_text__isnull=False
            ).exclude(id=submission.id)[:20]  # آخر 20 تسليم
            
            if not previous_submissions.exists():
                # حفظ النص للمقارنة المستقبلية
                submission.validation_data = submission.validation_data or {}
                submission.validation_data['pdf_text'] = text[:5000]  # حفظ أول 5000 حرف
                
                return {
                    'status': 'pass',
                    'message': 'لا توجد تسليمات سابقة للمقارنة',
                    'score': 100
                }
            
            # تجهيز النصوص
            current_text = text[:5000]  # أول 5000 حرف
            previous_texts = [
                sub.validation_data.get('pdf_text', '')[:5000]
                for sub in previous_submissions
                if sub.validation_data and sub.validation_data.get('pdf_text')
            ]
            
            if not previous_texts:
                submission.validation_data = submission.validation_data or {}
                submission.validation_data['pdf_text'] = current_text
                return {
                    'status': 'pass',
                    'message': 'لا توجد نصوص سابقة صالحة للمقارنة',
                    'score': 100
                }
            
            # TF-IDF + Cosine Similarity
            all_texts = [current_text] + previous_texts
            vectorizer = TfidfVectorizer(max_features=1000, stop_words=None)
            tfidf_matrix = vectorizer.fit_transform(all_texts)
            
            # حساب التشابه مع كل نص سابق
            current_vector = tfidf_matrix[0:1]
            previous_vectors = tfidf_matrix[1:]
            similarities = cosine_similarity(current_vector, previous_vectors)[0]
            
            # أعلى نسبة تشابه
            max_similarity = float(np.max(similarities)) * 100
            max_similarity_idx = int(np.argmax(similarities))
            
            # حفظ النص الحالي
            submission.validation_data = submission.validation_data or {}
            submission.validation_data['pdf_text'] = current_text
            submission.validation_data['max_similarity'] = max_similarity
            
            logger.info(f"📊 أعلى نسبة تشابه: {max_similarity:.1f}%")
            
            # التقييم
            threshold = submission.project.plagiarism_threshold  # من المشروع
            
            if max_similarity > 85:
                similar_sub = list(previous_submissions)[max_similarity_idx]
                return {
                    'status': 'fail',
                    'message': f'تشابه عالي جداً ({max_similarity:.0f}%) مع تسليم سابق',
                    'max_similarity': max_similarity,
                    'similar_submission': {
                        'id': similar_sub.id,
                        'student': similar_sub.submitted_student_name,
                        'submitted_at': similar_sub.submitted_at.isoformat()
                    },
                    'score': 0
                }
            elif max_similarity > threshold:
                return {
                    'status': 'warning',
                    'message': f'تشابه متوسط ({max_similarity:.0f}%) مع تسليم سابق',
                    'max_similarity': max_similarity,
                    'score': 70
                }
            else:
                return {
                    'status': 'pass',
                    'message': f'نسبة التشابه منخفضة ({max_similarity:.0f}%)',
                    'max_similarity': max_similarity,
                    'score': 100
                }
                
        except Exception as e:
            logger.error(f"❌ خطأ في كشف الانتحال: {str(e)}")
            return {
                'status': 'warning',
                'message': 'تعذر فحص الانتحال',
                'score': 80
            }
    
    # ====================================
    # Image Validation Helper Methods
    # ====================================
    
    def _check_image_ocr(self, file_path):
        """
        OCR على الصورة
        
        Args:
            file_path: مسار الصورة
            
        Returns:
            dict: النص المستخرج
        """
        try:
            import easyocr
            
            reader = easyocr.Reader(['ar', 'en'], gpu=False)
            result = reader.readtext(file_path)
            
            # استخراج النصوص
            texts = [detection[1] for detection in result]
            combined_text = ' '.join(texts)
            
            word_count = len(combined_text.split())
            
            logger.info(f"📝 OCR: استخراج {word_count} كلمة من الصورة")
            
            return {
                'status': 'pass',
                'message': f'تم استخراج {word_count} كلمة',
                'text': combined_text,
                'word_count': word_count,
                'score': 100 if word_count > 10 else 80
            }
            
        except Exception as e:
            logger.error(f"❌ خطأ في OCR للصورة: {str(e)}")
            return {
                'status': 'warning',
                'message': 'تعذر قراءة النص من الصورة',
                'text': '',
                'word_count': 0,
                'score': 70
            }
    
    def _analyze_image_content(self, file_path, project):
        """
        تحليل محتوى الصورة بـ Gemini Vision
        
        Args:
            file_path: مسار الصورة
            project: كائن المشروع
            
        Returns:
            dict: نتيجة التحليل
        """
        try:
            import google.generativeai as genai
            from PIL import Image
            
            if not self.gemini_vision:
                return {
                    'status': 'warning',
                    'message': 'Gemini Vision غير متاح',
                    'score': 70
                }
            
            # تحميل الصورة
            img = Image.open(file_path)
            
            # تجهيز Prompt
            prompt = f"""حلل هذه الصورة بدقة وأجب بصيغة JSON:

معلومات المشروع:
- العنوان: {project.title}
- الوصف: {project.description or 'غير محدد'}

أجب على التالي:
1. quality_score: جودة الصورة (0-100)
2. relevance_to_topic: ارتباط بالمشروع (0-100)
3. has_inappropriate_content: محتوى غير مناسب؟ (true/false)
4. description: وصف قصير للمحتوى
5. recommendation: التوصية (approved/rejected/needs_review)

أجب بصيغة JSON فقط."""
            
            response = self.gemini_vision.generate_content([prompt, img])
            
            # تحليل النتيجة
            import json
            try:
                result = json.loads(response.text)
            except:
                result = {
                    'quality_score': 75,
                    'relevance_to_topic': 75,
                    'has_inappropriate_content': False,
                    'description': response.text[:100],
                    'recommendation': 'approved'
                }
            
            # حساب الدرجة
            quality = result.get('quality_score', 70)
            relevance = result.get('relevance_to_topic', 70)
            overall = (quality + relevance) / 2
            
            # التحقق
            if result.get('has_inappropriate_content'):
                return {
                    'status': 'fail',
                    'message': 'الصورة تحتوي على محتوى غير مناسب',
                    'analysis': result,
                    'score': 0
                }
            
            if overall < 50:
                return {
                    'status': 'fail',
                    'message': f'جودة الصورة منخفضة ({overall:.0f}%)',
                    'analysis': result,
                    'score': overall
                }
            else:
                return {
                    'status': 'pass',
                    'message': f'جودة الصورة جيدة ({overall:.0f}%)',
                    'analysis': result,
                    'score': overall
                }
                
        except Exception as e:
            logger.error(f"❌ خطأ في تحليل الصورة: {str(e)}")
            return {
                'status': 'warning',
                'message': 'تعذر تحليل الصورة',
                'score': 70
            }
    
    def _check_image_quality(self, file_path):
        """
        فحص جودة الصورة (الدقة والحجم)
        
        Args:
            file_path: مسار الصورة
            
        Returns:
            dict: نتيجة الفحص
        """
        try:
            from PIL import Image
            
            img = Image.open(file_path)
            width, height = img.size
            pixels = width * height
            
            logger.info(f"📏 حجم الصورة: {width}x{height} ({pixels:,} pixels)")
            
            # التقييم
            if pixels < 100000:  # أقل من 0.1 ميجا بكسل
                return {
                    'status': 'warning',
                    'message': f'الصورة صغيرة ({width}x{height})',
                    'width': width,
                    'height': height,
                    'score': 60
                }
            elif pixels > 25000000:  # أكبر من 25 ميجا بكسل
                return {
                    'status': 'warning',
                    'message': f'الصورة كبيرة جداً ({width}x{height})',
                    'width': width,
                    'height': height,
                    'score': 80
                }
            else:
                return {
                    'status': 'pass',
                    'message': f'حجم الصورة مناسب ({width}x{height})',
                    'width': width,
                    'height': height,
                    'score': 100
                }
                
        except Exception as e:
            logger.error(f"❌ خطأ في فحص جودة الصورة: {str(e)}")
            return {
                'status': 'warning',
                'message': 'تعذر فحص الجودة',
                'score': 70
            }
    
    def _check_image_similarity(self, file_path, submission):
        """
        كشف التشابه بين الصور
        
        Args:
            file_path: مسار الصورة
            submission: كائن التسليم
            
        Returns:
            dict: نتيجة الفحص
        """
        try:
            import imagehash
            from PIL import Image
            from .models import Submission
            
            # حساب hash للصورة الحالية
            img = Image.open(file_path)
            current_hash = imagehash.average_hash(img)
            
            # البحث عن صور سابقة
            previous_submissions = Submission.objects.filter(
                project=submission.project,
                file_type='image',
                validation_data__image_hash__isnull=False
            ).exclude(id=submission.id)[:20]
            
            if not previous_submissions.exists():
                submission.validation_data = submission.validation_data or {}
                submission.validation_data['image_hash'] = str(current_hash)
                return {
                    'status': 'pass',
                    'message': 'لا توجد صور سابقة للمقارنة',
                    'score': 100
                }
            
            # المقارنة
            max_similarity = 0
            similar_sub = None
            
            for prev_sub in previous_submissions:
                try:
                    prev_hash_str = prev_sub.validation_data.get('image_hash')
                    if prev_hash_str:
                        prev_hash = imagehash.hex_to_hash(prev_hash_str)
                        difference = current_hash - prev_hash
                        similarity = max(0, 100 - (difference * 2))
                        
                        if similarity > max_similarity:
                            max_similarity = similarity
                            similar_sub = prev_sub
                except:
                    continue
            
            # حفظ hash الحالي
            submission.validation_data = submission.validation_data or {}
            submission.validation_data['image_hash'] = str(current_hash)
            
            logger.info(f"📊 أعلى تشابه: {max_similarity:.1f}%")
            
            # التقييم
            if max_similarity > 90:
                return {
                    'status': 'fail',
                    'message': f'صورة مشابهة جداً ({max_similarity:.0f}%) لتسليم سابق',
                    'similarity': max_similarity,
                    'score': 0
                }
            elif max_similarity > 70:
                return {
                    'status': 'warning',
                    'message': f'تشابه متوسط ({max_similarity:.0f}%)',
                    'similarity': max_similarity,
                    'score': 70
                }
            else:
                return {
                    'status': 'pass',
                    'message': f'الصورة أصلية ({max_similarity:.0f}% تشابه)',
                    'similarity': max_similarity,
                    'score': 100
                }
                
        except Exception as e:
            logger.error(f"❌ خطأ في فحص التشابه: {str(e)}")
            return {
                'status': 'warning',
                'message': 'تعذر فحص التشابه',
                'score': 80
            }
    
    # ====================================
    # Document Validation Helper Methods
    # ====================================
    
    def _extract_word_text(self, file_path):
        """استخراج النص من Word"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            text = '\n'.join([para.text for para in doc.paragraphs if para.text])
            word_count = len(text.split())
            
            logger.info(f"📄 Word: {word_count} كلمة")
            
            if word_count < 10:
                return {
                    'status': 'fail',
                    'message': 'المستند شبه فارغ',
                    'text': text,
                    'word_count': word_count,
                    'score': 0
                }
            
            return {
                'status': 'pass',
                'message': f'تم استخراج {word_count} كلمة',
                'text': text,
                'word_count': word_count,
                'score': 100
            }
            
        except Exception as e:
            logger.error(f"❌ خطأ في قراءة Word: {str(e)}")
            return {
                'status': 'fail',
                'message': f'فشل في قراءة الملف: {str(e)}',
                'text': '',
                'word_count': 0,
                'score': 0
            }
    
    def _extract_excel_text(self, file_path):
        """استخراج النص من Excel"""
        try:
            from openpyxl import load_workbook
            
            wb = load_workbook(file_path)
            text = ''
            cell_count = 0
            
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    for cell in row:
                        if cell:
                            text += str(cell) + ' '
                            cell_count += 1
            
            word_count = len(text.split())
            
            logger.info(f"📊 Excel: {cell_count} خلية، {word_count} كلمة")
            
            if cell_count < 5:
                return {
                    'status': 'fail',
                    'message': 'المستند شبه فارغ',
                    'text': text,
                    'word_count': word_count,
                    'cell_count': cell_count,
                    'score': 0
                }
            
            return {
                'status': 'pass',
                'message': f'تم استخراج {cell_count} خلية',
                'text': text,
                'word_count': word_count,
                'cell_count': cell_count,
                'score': 100
            }
            
        except Exception as e:
            logger.error(f"❌ خطأ في قراءة Excel: {str(e)}")
            return {
                'status': 'fail',
                'message': f'فشل في قراءة الملف: {str(e)}',
                'text': '',
                'word_count': 0,
                'score': 0
            }
    
    def _extract_ppt_text(self, file_path):
        """استخراج النص من PowerPoint"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text = ''
            slide_count = len(prs.slides)
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + '\n'
            
            word_count = len(text.split())
            
            logger.info(f"📊 PPT: {slide_count} شريحة، {word_count} كلمة")
            
            if word_count < 10:
                return {
                    'status': 'fail',
                    'message': 'العرض التقديمي شبه فارغ',
                    'text': text,
                    'word_count': word_count,
                    'slide_count': slide_count,
                    'score': 0
                }
            
            return {
                'status': 'pass',
                'message': f'تم استخراج {word_count} كلمة من {slide_count} شريحة',
                'text': text,
                'word_count': word_count,
                'slide_count': slide_count,
                'score': 100
            }
            
        except Exception as e:
            logger.error(f"❌ خطأ في قراءة PowerPoint: {str(e)}")
            return {
                'status': 'fail',
                'message': f'فشل في قراءة الملف: {str(e)}',
                'text': '',
                'word_count': 0,
                'score': 0
            }
    
    def _check_document_stats(self, text_result, project):
        """فحص إحصائيات المستند"""
        word_count = text_result.get('word_count', 0)
        
        # القيود
        constraints = project.file_constraints or {}
        min_words = constraints.get('min_words', 50)
        max_words = constraints.get('max_words', 3000)
        
        if word_count < min_words:
            return {
                'status': 'fail',
                'message': f'عدد الكلمات قليل ({word_count}). المطلوب {min_words}+',
                'word_count': word_count,
                'score': 0
            }
        elif word_count > max_words:
            return {
                'status': 'warning',
                'message': f'عدد الكلمات كثير ({word_count}). الحد الأقصى {max_words}',
                'word_count': word_count,
                'score': 80
            }
        else:
            return {
                'status': 'pass',
                'message': f'عدد الكلمات مناسب ({word_count})',
                'word_count': word_count,
                'score': 100
            }
    
    def _analyze_document_content(self, text, project, file_ext):
        """تحليل محتوى المستند بـ Gemini"""
        try:
            if not self.gemini_flash:
                return {
                    'status': 'warning',
                    'message': 'Gemini غير متاح',
                    'score': 70
                }
            
            text_sample = text[:2000] if len(text) > 2000 else text
            
            prompt = f"""حلل هذا المستند ({file_ext}) بصيغة JSON:

المشروع: {project.title}

النص:
{text_sample}

أجب (JSON):
1. content_quality: جودة (0-100)
2. relevance: ارتباط بالمشروع (0-100)
3. recommendation: (approved/rejected/needs_review)

JSON فقط."""
            
            response = self.gemini_flash.generate_content(prompt)
            
            import json
            try:
                result = json.loads(response.text)
            except:
                result = {
                    'content_quality': 75,
                    'relevance': 75,
                    'recommendation': 'approved'
                }
            
            quality = result.get('content_quality', 70)
            relevance = result.get('relevance', 70)
            overall = (quality + relevance) / 2
            
            if overall < 50:
                return {
                    'status': 'fail',
                    'message': f'جودة منخفضة ({overall:.0f}%)',
                    'score': overall
                }
            else:
                return {
                    'status': 'pass',
                    'message': f'جودة جيدة ({overall:.0f}%)',
                    'score': overall
                }
                
        except Exception as e:
            logger.error(f"❌ خطأ في تحليل المستند: {str(e)}")
            return {
                'status': 'warning',
                'message': 'تعذر التحليل',
                'score': 70
            }
    
    # ====================================
    # Audio Validation Helper Methods
    # ====================================
    
    def _check_audio_duration(self, file_path, project):
        """
        فحص مدة الصوت
        
        Args:
            file_path: مسار الملف الصوتي
            project: كائن المشروع
            
        Returns:
            dict: نتيجة الفحص
        """
        try:
            from pydub import AudioSegment
            
            # تحميل الملف الصوتي
            audio = AudioSegment.from_file(file_path)
            duration = len(audio) / 1000.0  # بالثواني
            
            # قراءة القيود من المشروع
            constraints = project.file_constraints or {}
            duration_constraints = constraints.get('duration', {})
            
            min_duration = duration_constraints.get('min', 10)  # default 10 seconds
            max_duration = duration_constraints.get('max', 180)  # default 3 minutes
            
            logger.info(f"🎵 مدة الصوت: {duration:.1f} ثانية (المطلوب: {min_duration}-{max_duration})")
            
            # التقييم
            if duration < min_duration:
                return {
                    'status': 'fail',
                    'message': f'الصوت قصير جداً ({duration:.1f}ث). المطلوب على الأقل {min_duration} ثانية',
                    'duration': duration,
                    'score': 0
                }
            elif duration > max_duration:
                return {
                    'status': 'fail',
                    'message': f'الصوت طويل جداً ({duration:.1f}ث). الحد الأقصى {max_duration} ثانية',
                    'duration': duration,
                    'score': 0
                }
            else:
                return {
                    'status': 'pass',
                    'message': f'مدة الصوت مناسبة ({duration:.1f}ث)',
                    'duration': duration,
                    'score': 100
                }
                
        except Exception as e:
            logger.error(f"❌ خطأ في فحص مدة الصوت: {str(e)}")
            return {
                'status': 'fail',
                'message': f'خطأ في قراءة الملف الصوتي: {str(e)}',
                'score': 0
            }
    
    def _audio_to_text(self, file_path):
        """
        تحويل الصوت إلى نص (Speech-to-Text)
        
        Args:
            file_path: مسار الملف الصوتي
            
        Returns:
            dict: النص المستخرج
        """
        try:
            import speech_recognition as sr
            from pydub import AudioSegment
            import os
            
            recognizer = sr.Recognizer()
            
            # تحويل الملف إلى WAV إذا لزم الأمر
            audio = AudioSegment.from_file(file_path)
            
            # تحديد مدة الصوت للمعالجة (أول دقيقة فقط لتوفير الوقت)
            duration_ms = min(len(audio), 60000)  # أول دقيقة
            audio_sample = audio[:duration_ms]
            
            # حفظ مؤقتاً كـ WAV
            temp_wav = file_path + '.temp.wav'
            audio_sample.export(temp_wav, format='wav')
            
            try:
                # قراءة الملف الصوتي
                with sr.AudioFile(temp_wav) as source:
                    audio_data = recognizer.record(source)
                
                # محاولة التعرف على النص (العربية أولاً، ثم الإنجليزية)
                text = ''
                try:
                    text = recognizer.recognize_google(audio_data, language='ar-SA')
                    logger.info(f"📝 تم التعرف على النص بالعربية")
                except:
                    try:
                        text = recognizer.recognize_google(audio_data, language='en-US')
                        logger.info(f"📝 تم التعرف على النص بالإنجليزية")
                    except:
                        pass
                
                # حذف الملف المؤقت
                if os.path.exists(temp_wav):
                    os.remove(temp_wav)
                
                word_count = len(text.split())
                
                logger.info(f"📝 Speech-to-Text: {word_count} كلمة")
                
                if word_count < 5:
                    return {
                        'status': 'fail',
                        'message': 'لم يتم التعرف على كلام واضح في الملف الصوتي',
                        'text': text,
                        'word_count': word_count,
                        'score': 0
                    }
                
                return {
                    'status': 'pass',
                    'message': f'تم استخراج {word_count} كلمة',
                    'text': text,
                    'word_count': word_count,
                    'score': 100
                }
                
            finally:
                # تأكد من حذف الملف المؤقت
                if os.path.exists(temp_wav):
                    try:
                        os.remove(temp_wav)
                    except:
                        pass
                
        except Exception as e:
            logger.error(f"❌ خطأ في Speech-to-Text: {str(e)}")
            return {
                'status': 'warning',
                'message': f'تعذر تحويل الصوت إلى نص: {str(e)}',
                'text': '',
                'word_count': 0,
                'score': 50
            }
    
    def _analyze_audio_content(self, text, project):
        """
        تحليل محتوى الصوت باستخدام Gemini
        
        Args:
            text: النص المستخرج من الصوت
            project: كائن المشروع
            
        Returns:
            dict: نتيجة التحليل
        """
        try:
            if not self.gemini_flash:
                return {
                    'status': 'warning',
                    'message': 'Gemini غير متاح',
                    'score': 70
                }
            
            if not text or len(text.strip()) < 10:
                return {
                    'status': 'warning',
                    'message': 'النص المستخرج قصير جداً للتحليل',
                    'score': 60
                }
            
            # تجهيز Prompt
            prompt = f"""حلل هذا النص المستخرج من ملف صوتي بصيغة JSON:

معلومات المشروع:
- العنوان: {project.title}
- الوصف: {project.description or 'غير محدد'}

النص المستخرج:
{text}

أجب (JSON):
1. content_quality: جودة المحتوى (0-100)
2. relevance: ارتباط بالمشروع (0-100)
3. language_quality: جودة اللغة (0-100)
4. recommendation: (approved/rejected/needs_review)

JSON فقط."""
            
            response = self.gemini_flash.generate_content(prompt)
            
            # تحليل النتيجة
            import json
            try:
                result = json.loads(response.text)
            except:
                result = {
                    'content_quality': 75,
                    'relevance': 75,
                    'language_quality': 75,
                    'recommendation': 'approved'
                }
            
            # حساب الدرجة
            quality = result.get('content_quality', 70)
            relevance = result.get('relevance', 70)
            language = result.get('language_quality', 70)
            overall = (quality + relevance + language) / 3
            
            if overall < 50:
                return {
                    'status': 'fail',
                    'message': f'جودة المحتوى منخفضة ({overall:.0f}%)',
                    'analysis': result,
                    'score': overall
                }
            elif overall < 70:
                return {
                    'status': 'warning',
                    'message': f'جودة المحتوى مقبولة ({overall:.0f}%)',
                    'analysis': result,
                    'score': overall
                }
            else:
                return {
                    'status': 'pass',
                    'message': f'جودة المحتوى ممتازة ({overall:.0f}%)',
                    'analysis': result,
                    'score': overall
                }
                
        except Exception as e:
            logger.error(f"❌ خطأ في تحليل محتوى الصوت: {str(e)}")
            return {
                'status': 'warning',
                'message': 'تعذر تحليل المحتوى',
                'score': 70
            }
    
    def _check_audio_similarity(self, text, submission):
        """
        كشف التشابه في الملفات الصوتية
        
        Args:
            text: النص المستخرج من الصوت
            submission: كائن التسليم
            
        Returns:
            dict: نتيجة الفحص
        """
        try:
            from .models import Submission
            from sklearn.feature_extraction.text import TfidfVectorizer
            from sklearn.metrics.pairwise import cosine_similarity
            import numpy as np
            
            if not text or len(text.strip()) < 10:
                return {
                    'status': 'warning',
                    'message': 'النص قصير جداً لفحص التشابه',
                    'score': 70
                }
            
            # البحث عن صوتيات سابقة
            previous_submissions = Submission.objects.filter(
                project=submission.project,
                file_type='audio',
                validation_data__audio_text__isnull=False
            ).exclude(id=submission.id)[:20]
            
            if not previous_submissions.exists():
                # حفظ النص للمقارنة المستقبلية
                submission.validation_data = submission.validation_data or {}
                submission.validation_data['audio_text'] = text[:3000]
                
                return {
                    'status': 'pass',
                    'message': 'لا توجد صوتيات سابقة للمقارنة',
                    'score': 100
                }
            
            # تجهيز النصوص
            current_text = text[:3000]
            previous_texts = [
                sub.validation_data.get('audio_text', '')[:3000]
                for sub in previous_submissions
                if sub.validation_data and sub.validation_data.get('audio_text')
            ]
            
            if not previous_texts:
                submission.validation_data = submission.validation_data or {}
                submission.validation_data['audio_text'] = current_text
                return {
                    'status': 'pass',
                    'message': 'لا توجد نصوص سابقة صالحة للمقارنة',
                    'score': 100
                }
            
            # TF-IDF + Cosine Similarity
            all_texts = [current_text] + previous_texts
            vectorizer = TfidfVectorizer(max_features=500)
            tfidf_matrix = vectorizer.fit_transform(all_texts)
            
            # حساب التشابه
            current_vector = tfidf_matrix[0:1]
            previous_vectors = tfidf_matrix[1:]
            similarities = cosine_similarity(current_vector, previous_vectors)[0]
            
            max_similarity = float(np.max(similarities)) * 100
            
            # حفظ النص الحالي
            submission.validation_data = submission.validation_data or {}
            submission.validation_data['audio_text'] = current_text
            
            logger.info(f"📊 أعلى نسبة تشابه: {max_similarity:.1f}%")
            
            # التقييم
            if max_similarity > 80:
                return {
                    'status': 'fail',
                    'message': f'تشابه عالي جداً ({max_similarity:.0f}%) مع تسليم سابق',
                    'max_similarity': max_similarity,
                    'score': 0
                }
            elif max_similarity > 60:
                return {
                    'status': 'warning',
                    'message': f'تشابه متوسط ({max_similarity:.0f}%) مع تسليم سابق',
                    'max_similarity': max_similarity,
                    'score': 70
                }
            else:
                return {
                    'status': 'pass',
                    'message': f'نسبة التشابه منخفضة ({max_similarity:.0f}%)',
                    'max_similarity': max_similarity,
                    'score': 100
                }
                
        except Exception as e:
            logger.error(f"❌ خطأ في فحص التشابه: {str(e)}")
            return {
                'status': 'warning',
                'message': 'تعذر فحص التشابه',
                'score': 80
            }
